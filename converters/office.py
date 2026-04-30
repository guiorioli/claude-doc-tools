from __future__ import annotations

from pathlib import Path
from typing import List


def convert_docx(path: Path, image_dir: Path) -> tuple[str, List[Path]]:
    import mammoth

    saved: List[Path] = []
    counter = {"i": 0}

    def handle_image(image):
        counter["i"] += 1
        ext = (image.content_type or "image/png").split("/")[-1]
        if ext == "jpeg":
            ext = "jpg"
        out = image_dir / f"image{counter['i']}.{ext}"
        with image.open() as src, open(out, "wb") as dst:
            dst.write(src.read())
        saved.append(out)
        return {"src": str(out)}

    with open(path, "rb") as fh:
        result = mammoth.convert_to_markdown(
            fh, convert_image=mammoth.images.img_element(handle_image)
        )
    return result.value, saved


def convert_xlsx(path: Path, image_dir: Path) -> tuple[str, List[Path]]:
    return _convert_excel(path, engine="openpyxl")


def convert_xls(path: Path, image_dir: Path) -> tuple[str, List[Path]]:
    return _convert_excel(path, engine="xlrd")


def _convert_excel(path: Path, engine: str) -> tuple[str, List[Path]]:
    import pandas as pd

    sheets = pd.read_excel(path, sheet_name=None, engine=engine)
    parts: List[str] = []
    for name, df in sheets.items():
        parts.append(f"## Sheet: {name}\n")
        if df.empty:
            parts.append("_(empty)_\n")
        else:
            parts.append(df.to_markdown(index=False))
            parts.append("")
    return "\n".join(parts), []


def convert_pptx(path: Path, image_dir: Path) -> tuple[str, List[Path]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    saved: List[Path] = []
    counter = 0
    parts: List[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {idx}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
                    parts.append("")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                counter += 1
                image = shape.image
                ext = image.ext or "png"
                out = image_dir / f"slide{idx}_image{counter}.{ext}"
                out.write_bytes(image.blob)
                saved.append(out)
                parts.append(f"![image]({out})")
                parts.append("")

        notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        if notes:
            parts.append("**Speaker notes:**")
            parts.append(notes)
            parts.append("")

    return "\n".join(parts), saved
